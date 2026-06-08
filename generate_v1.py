import streamlit as st
import pandas as pd
import os
from io import BytesIO
import plotly.express as px
from pydantic import BaseModel, Field
from typing import List, Optional
from openai import AzureOpenAI
import sqlite3
import uuid
from datetime import datetime
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

api_key = os.getenv("AZURE_OPENAI_API_KEY")
azure_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
azure_deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o-mini")

# PPTX imports
from pptx import Presentation
from pptx.util import Inches

# LangGraph / LangChain imports
from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, AIMessage
from langchain_core.tools import tool
from langgraph.prebuilt import create_react_agent

# ==========================================
# 1. Structured Data Models (Pydantic)
# ==========================================
class InsightBullet(BaseModel):
    metric_impacted: str = Field(description="The specific metric or KPI this insight refers to.")
    observation: str = Field(description="A clear statement of what happened based on the data facts provided.")
    business_implication: str = Field(description="Why this matters to stakeholders or operations.")

class AutomatedCommentary(BaseModel):
    headline: str = Field(description="A punchy, informative macro-level headline summarizing the performance.")
    executive_summary: str = Field(description="A brief 2-3 sentence overview of the data trend.")
    key_insights: List[InsightBullet] = Field(description="List of specific mathematical insights extracted from the data.")
    recommended_actions: List[str] = Field(description="Actionable next steps based on these findings.")

# ==========================================
# 2. Roles & Definitions
# ==========================================
ROLES = {
    "CEO / Executive": "Focus on high-level ROI, overall growth, strategic risks, and enterprise value. Use executive, concise language.",
    "Sales Director": "Focus on pipeline velocity, conversion rates, regional performance, and revenue targets. Use sales-oriented, aggressive language.",
    "Operations Lead": "Focus on cost-efficiency, bottlenecks, process improvements, and resource allocation. Use practical, metric-driven language."
}

# ==========================================
# 3. LangGraph Tool Definition
# ==========================================
@tool
def analyze_dataframe(query: str) -> str:
    """
    Executes Python pandas code on the uploaded dataframe 'df'.
    Example query: "df['Revenue'].sum()" or "df.groupby('Region')['CAC'].mean()"
    """
    if "df" not in st.session_state or st.session_state.df is None:
        return "Error: No dataframe available."
    
    df = st.session_state.df
    try:
        # Evaluate the pandas operation
        result = eval(query)
        return str(result)
    except Exception as e:
        return f"Error executing code '{query}': {str(e)}"

# ==========================================
# 4. Helpers: Auto-Charts & PPTX Generation
# ==========================================
def generate_auto_charts(df: pd.DataFrame) -> list:
    """Intelligently guesses and builds interactive Plotly charts based on column types."""
    charts = []
    num_cols = df.select_dtypes(include='number').columns.tolist()
    cat_cols = df.select_dtypes(exclude='number').columns.tolist()
    
    if cat_cols and num_cols:
        # 1. Bar Chart (Categorical vs Numeric)
        fig1 = px.bar(df, x=cat_cols[0], y=num_cols[0], title=f"{num_cols[0]} segmented by {cat_cols[0]}", template="plotly_dark" if st.get_option("theme.base") == "dark" else "plotly_white")
        fig1.update_layout(margin=dict(l=20, r=20, t=40, b=20))
        charts.append({"name": f"Bar Analysis: {cat_cols[0]}", "fig": fig1})
        
        # 2. Line Chart / Trend (If multiple numerics exist)
        if len(num_cols) > 1:
            fig2 = px.line(df, x=cat_cols[0], y=num_cols[1], title=f"{num_cols[1]} Trend over {cat_cols[0]}", markers=True, template="plotly_dark" if st.get_option("theme.base") == "dark" else "plotly_white")
            fig2.update_layout(margin=dict(l=20, r=20, t=40, b=20))
            charts.append({"name": f"Trend Analysis: {num_cols[1]}", "fig": fig2})
            
        # 3. Scatter Plot (Numeric vs Numeric)
        if len(num_cols) >= 2:
            fig3 = px.scatter(df, x=num_cols[0], y=num_cols[1], color=cat_cols[0] if cat_cols else None, title=f"Correlation: {num_cols[0]} vs {num_cols[1]}", template="plotly_dark" if st.get_option("theme.base") == "dark" else "plotly_white")
            charts.append({"name": f"Scatter Analysis", "fig": fig3})
            
    return charts

def create_ppt(commentary: AutomatedCommentary, role: str, custom_headline: str, custom_summary: str, charts: list) -> BytesIO:
    """Generates a presentation dynamically inserting approved charts and edited narrative."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    chart_layout = prs.slide_layouts[5] # Blank/Title only layout

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = custom_headline
    slide.placeholders[1].text = f"Automated Data Briefing\nPrepared for: {role}"

    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Executive Summary"
    slide.placeholders[1].text = custom_summary

    # Slide 3: Key Insights
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Key Insights"
    tf = slide.placeholders[1].text_frame
    tf.text = "Data Observations:"
    for insight in commentary.key_insights:
        p = tf.add_paragraph()
        p.text = f"• {insight.metric_impacted}: {insight.observation} ({insight.business_implication})"
        p.level = 1

    # Slide 4: Action Items
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Recommended Actions"
    tf = slide.placeholders[1].text_frame
    for action in commentary.recommended_actions:
        p = tf.add_paragraph()
        p.text = action
        p.level = 0

    # Slide 5+: Selected Interactive Charts
    for fig_data in charts:
        slide = prs.slides.add_slide(chart_layout)
        slide.shapes.title.text = fig_data["name"]
        
        # Convert Plotly to static PNG for PPTX
        image_stream = BytesIO(fig_data["fig"].to_image(format="png", width=800, height=450))
        slide.shapes.add_picture(image_stream, Inches(1), Inches(2), width=Inches(8))

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# ==========================================
# 4.5 SQLite Database & Memory Management
# ==========================================
def init_db():
    """Initializes the SQLite database for chat persistence."""
    conn = sqlite3.connect("chat_history.db")
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS messages
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  session_id TEXT,
                  role TEXT,
                  content TEXT,
                  timestamp DATETIME DEFAULT CURRENT_TIMESTAMP)''')
    conn.commit()
    conn.close()

def save_message(session_id: str, role: str, content: str):
    """Saves a single chat message to the database."""
    conn = sqlite3.connect("chat_history.db")
    c = conn.cursor()
    c.execute("INSERT INTO messages (session_id, role, content) VALUES (?, ?, ?)",
              (session_id, role, content))
    conn.commit()
    conn.close()

def load_chat_history(session_id: str) -> list:
    """Loads previous messages for a specific session."""
    conn = sqlite3.connect("chat_history.db")
    c = conn.cursor()
    c.execute("SELECT role, content FROM messages WHERE session_id = ? ORDER BY timestamp ASC", (session_id,))
    rows = c.fetchall()
    conn.close()
    return [{"role": r[0], "content": r[1]} for r in rows]

def get_all_sessions() -> list:
    """Retrieves all unique session IDs sorted by newest first."""
    conn = sqlite3.connect("chat_history.db")
    c = conn.cursor()
    c.execute("SELECT DISTINCT session_id FROM messages ORDER BY timestamp DESC")
    rows = c.fetchall()
    conn.close()
    return [r[0] for r in rows]

# Initialize Database on startup
init_db()

# ==========================================
# 5. Main Streamlit UI Setup
# ==========================================
st.set_page_config(page_title="Dynamic Data Storyteller", layout="wide", page_icon="📊")

# Initialize Session State Variables
for state_var in ["messages", "commentary", "summary_stats", "df", "generated_charts", "role"]:
    if state_var not in st.session_state:
        st.session_state[state_var] = [] if state_var in ["messages", "generated_charts"] else None

# Initialize Unique Session ID for Chat Memory
if "session_id" not in st.session_state:
    st.session_state.session_id = str(uuid.uuid4())

# Sidebar Configuration
with st.sidebar:
    st.header("⚙️ Configuration")
    
    selected_role = st.selectbox("Target Audience Role", list(ROLES.keys()))
    context_notes = st.text_area("Additional Business Context", placeholder="E.g., Q2 experienced an aggressive ad campaign...")
    st.session_state.role = selected_role
    
    st.divider()
    st.header("🗄️ Chat Memory")
    
    # Button to start a completely fresh conversation memory
    if st.button("➕ Start New Chat Session", use_container_width=True):
        st.session_state.session_id = str(uuid.uuid4())
        st.session_state.messages = []
        st.rerun()

    # Dropdown to load previous sessions from SQLite
    sessions = get_all_sessions()
    if sessions:
        selected_session = st.selectbox("Load Previous Session", ["Current"] + sessions)
        if selected_session != "Current" and selected_session != st.session_state.session_id:
            st.session_state.session_id = selected_session
            st.session_state.messages = load_chat_history(selected_session)

st.title("📊 Dynamic Data Storyteller")
st.markdown("Upload your CSV, review auto-generated insights, interact with your data, and export directly to PowerPoint.")

# Step 1: File Uploader
uploaded_file = st.file_uploader("Upload CSV or Excel Data", type=["csv", "xlsx"])

if uploaded_file:
    # Button to trigger analysis (prevents infinite reruns on minor edits)
    if st.button("🚀 Analyze & Generate Draft Commentary", type="primary", use_container_width=True):
        with st.spinner("Analyzing data structure & rendering narrative..."):
            
            # Load Data
            df = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
            st.session_state.df = df
            
            # Basic summary stats for prompt context
            summary_stats = {
                "total_rows": int(len(df)),
                "columns": list(df.columns),
                "numeric_summary": df.describe().to_dict(),
            }
            st.session_state.summary_stats = summary_stats
            
            client = AzureOpenAI(
                api_key=api_key,
                azure_endpoint=azure_endpoint,
                api_version=api_version
            )
            
            # Formulate the Prompt
            system_instruction = (
                f"You are an expert data analyst. Target Audience: {selected_role}. "
                f"Tone Instructions: {ROLES[selected_role]} "
                "Convert the raw statistical summaries into actionable commentary. Do not hallucinate."
            )
            user_prompt = f"Data Facts Summary:\n{summary_stats}\n\nAdditional Business Context:\n{context_notes}"

            # Structured Output generation
            try:
                completion = client.beta.chat.completions.parse(
                    model=azure_deployment,
                    messages=[
                        {"role": "system", "content": system_instruction},
                        {"role": "user", "content": user_prompt}
                    ],
                    response_format=AutomatedCommentary,
                    temperature=0.2
                )
                
                st.session_state.commentary = completion.choices[0].message.parsed
                st.session_state.generated_charts = generate_auto_charts(st.session_state.df)
                
                # Reset chat and generate new session when new file is analyzed
                st.session_state.session_id = str(uuid.uuid4())
                st.session_state.messages = [] 
                
            except Exception as e:
                st.error(f"Error generating commentary: {e}")

    # ==========================================
    # Step 2: Review, Edit, & Export UI
    # ==========================================
    if st.session_state.commentary:
        commentary = st.session_state.commentary
        st.divider()
        st.header("Step 2: Review, Edit & Export")
        
        col_text, col_charts = st.columns([1, 1.2])
        approved_charts = []
        
        with col_charts:
            st.subheader("📈 Interactive Visualizations")
            for i, chart_data in enumerate(st.session_state.generated_charts):
                with st.container(border=True):
                    st.plotly_chart(chart_data["fig"], use_container_width=True, key=f"plot_{i}")
                    if st.checkbox(f"✅ Include '{chart_data['name']}' in PPTX", value=True, key=f"check_{i}"):
                        approved_charts.append(chart_data)

        with col_text:
            st.subheader("📝 Edit Narrative")
            edited_headline = st.text_input("Slide Headline", value=commentary.headline)
            edited_summary = st.text_area("Executive Summary", value=commentary.executive_summary, height=120)
            
            st.markdown("#### Key Insights (Preview)")
            for idx, insight in enumerate(commentary.key_insights):
                st.markdown(f"- **{insight.metric_impacted}**: {insight.observation}")
        
            st.markdown("#### Recommended Actions (Preview)")
            for action in commentary.recommended_actions:
                st.markdown(f"- {action}")

            # Export Button
            st.divider()
            ppt_file = create_ppt(commentary, st.session_state.role, edited_headline, edited_summary, approved_charts)
            st.download_button(
                label="📥 Download Approved PowerPoint",
                data=ppt_file,
                file_name=f"{st.session_state.role.replace(' ', '_')}_Briefing.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True
            )
        
        # ==========================================
        # Step 3: LangGraph Agent Chat
        # ==========================================
        st.divider()
        st.subheader("🤖 Data Analyst Agent")
        st.caption(f"Session ID: `{st.session_state.session_id}` | Ask questions about your data. The agent writes Python pandas code in the background to answer you.")

        # Display history
        for msg in st.session_state.messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        # Chat input
        if prompt := st.chat_input("E.g., What is the total revenue for North America?"):
            # Save User Input to Session State and SQLite DB
            st.session_state.messages.append({"role": "user", "content": prompt})
            save_message(st.session_state.session_id, "user", prompt)
            
            with st.chat_message("user"):
                st.markdown(prompt)

            with st.chat_message("assistant"):
                with st.spinner("Analyzing data with Python..."):
                    try:
                        llm = AzureChatOpenAI(
                            azure_deployment=azure_deployment,
                            api_version=api_version,
                            azure_endpoint=azure_endpoint,
                            api_key=api_key,
                            temperature=0.0
                        )
                        tools = [analyze_dataframe]
                        
                        sys_msg = (
                            "You are a helpful data analyst. Use the `analyze_dataframe` tool to execute pandas "
                            "code on the user's data to answer their questions. Only give the final answer."
                        )
                        agent = create_react_agent(llm, tools, state_modifier=sys_msg)
                        
                        # Convert persistent state to LangChain format (This acts as the agent's memory)
                        lc_messages = [HumanMessage(content=m["content"]) if m["role"] == "user" else AIMessage(content=m["content"]) for m in st.session_state.messages]
                        
                        # Invoke Agent
                        response = agent.invoke({"messages": lc_messages})
                        answer = response["messages"][-1].content
                        
                        st.markdown(answer)
                        
                        # Save Agent Response to Session State and SQLite DB
                        st.session_state.messages.append({"role": "assistant", "content": answer})
                        save_message(st.session_state.session_id, "assistant", answer)
                    except Exception as e:
                        st.error(f"Chat Error: {e}")
