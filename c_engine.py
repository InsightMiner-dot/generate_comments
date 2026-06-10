import os
import io
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import streamlit as st
from typing import TypedDict, Annotated, Dict, List
from pydantic import BaseModel, Field

from dotenv import load_dotenv
from langgraph.graph import StateGraph, START, END
from langgraph.graph.message import add_messages
from langgraph.checkpoint.memory import MemorySaver

from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_openai import AzureChatOpenAI

from pptx import Presentation
from pptx.util import Inches

# ==========================================
# 1. Configuration & App Initialization
# ==========================================
load_dotenv()

st.set_page_config(
    page_title="Enterprise Insights Engine", 
    page_icon="📊", 
    layout="wide"
)

# User's Dynamic Azure OpenAI Configuration
azure_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
azure_api_key = os.getenv("AZURE_OPENAI_KEY")
azure_deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT")
azure_api_version = os.getenv("AZURE_OPENAI_API_VERSION")
model_name = os.getenv("MODEL_NAME")

if not azure_api_key or not azure_endpoint:
    st.error("Missing Azure OpenAI credentials. Please verify your `.env` file configuration.")
    st.stop()

# Initialize LLM
llm = AzureChatOpenAI(
    azure_endpoint=azure_endpoint,
    api_key=azure_api_key,
    api_version=azure_api_version,
    deployment_name=azure_deployment,
    model=model_name,
    temperature=0.0,
    max_tokens=500,
)

# ==========================================
# 2. Pydantic Models for Output Restructuring
# ==========================================
class SlideContent(BaseModel):
    title: str = Field(description="The title of the slide")
    bullet_points: List[str] = Field(description="3-5 bullet points summarizing key insights for this slide")

class ChartSpecification(BaseModel):
    chart_title: str = Field(description="Title of the chart")
    x_column: str = Field(description="The exact name of the column to use for the X-axis")
    y_column: str = Field(description="The exact name of the column to use for the Y-axis")
    chart_type: str = Field(description="Type of chart to generate: 'bar', 'line', or 'scatter'")

class PresentationDeck(BaseModel):
    slides: List[SlideContent] = Field(description="The list of slides in the presentation")
    chart_suggestion: ChartSpecification = Field(description="The optimal data visualization based on the columns")

# ==========================================
# 3. Graph State Definition
# ==========================================
class PresentationState(TypedDict):
    messages: Annotated[list, add_messages]
    schema_info: str
    data_summary: str
    missing_params: List[str]
    presentation_params: Dict[str, str]
    dataframe_json: str  
    output_ready: bool

# ==========================================
# 4. Graph Architecture Nodes
# ==========================================
def ingest_and_summarize(state: PresentationState):
    schema = state["schema_info"]
    prompt = f"Analyze this dataset schema and provide a brief, engaging summary of what this data represents:\n{schema}"
    response = llm.invoke([SystemMessage(content=prompt)])
    
    msg = (f"### Data Ingestion Summary\n\n{response.content}\n\n"
           "Let's craft your presentation framework. What type of **commentary** style are you looking for "
           "(e.g., Executive Summary, Growth Deep Dive, Operational Metric Analysis)?")
    
    return {
        "data_summary": response.content,
        "messages": [AIMessage(content=msg)]
    }

def gather_parameters(state: PresentationState):
    missing = list(state.get("missing_params", []))
    params = dict(state.get("presentation_params", {}))
    last_message = state["messages"][-1].content
    
    if "commentary_type" in missing:
        params["commentary_type"] = last_message
        missing.remove("commentary_type")
        next_question = "Understood. Who is the target **persona or audience** for this deck?"
    
    elif "persona" in missing:
        params["persona"] = last_message
        missing.remove("persona")
        next_question = "Got it. What are the key specific **topics, anomalies, or metrics** we should emphasize?"
        
    elif "topics" in missing:
        params["topics"] = last_message
        missing.remove("topics")
        next_question = "Perfect. How many total **slides/pages** would you like inside this deck?"
        
    elif "pages" in missing:
        params["pages"] = last_message
        missing.remove("pages")
        next_question = "Excellent. Finally, what should be the **Main Title** of this PowerPoint presentation?"
        
    elif "title" in missing:
        params["title"] = last_message
        missing.remove("title")
        next_question = "Everything is gathered! Review your settings and reply with **'Approve'** to assemble the final presentation deck."
    else:
        next_question = "Reply with **'Approve'** to run the deck compiler generator."

    return {
        "presentation_params": params, 
        "missing_params": missing, 
        "messages": [AIMessage(content=next_question)]
    }

def generate_presentation_content(state: PresentationState):
    params = state["presentation_params"]
    schema = state["schema_info"]
    summary = state["data_summary"]
    
    prompt = f"""
    Create a structured slide deck layout and choose a data visualization plot configuration based on this context:
    Dataset Schema: {schema}
    Data Summary: {summary}
    Title: {params.get('title')}
    Commentary Type: {params.get('commentary_type')}
    Audience/Persona: {params.get('persona')}
    Key Topics: {params.get('topics')}
    Target Number of Slides: {params.get('pages')}
    
    Ensure your chart suggestion maps perfectly to the available columns listed in the context.
    """
    
    structured_llm = llm.with_structured_output(PresentationDeck)
    deck_data = structured_llm.invoke([SystemMessage(content=prompt)])
    
    df = pd.read_json(io.StringIO(state["dataframe_json"]))
    chart_spec = deck_data.chart_suggestion
    chart_bytes = io.BytesIO()
    has_plot = False
    
    try:
        plt.figure(figsize=(8, 4.5))
        if chart_spec.chart_type == 'bar':
            df.head(10).plot(kind='bar', x=chart_spec.x_column, y=chart_spec.y_column, ax=plt.gca(), color="#2c3e50")
        elif chart_spec.chart_type == 'scatter':
            df.plot(kind='scatter', x=chart_spec.x_column, y=chart_spec.y_column, ax=plt.gca(), color="#e74c3c")
        else:
            df.head(10).plot(kind='line', x=chart_spec.x_column, y=chart_spec.y_column, ax=plt.gca(), linewidth=2.5)
            
        plt.title(chart_spec.chart_title, fontsize=14, fontweight='bold')
        plt.xticks(rotation=45, ha='right')
        plt.grid(axis='y', linestyle='--', alpha=0.7)
        plt.tight_layout()
        plt.savefig(chart_bytes, format='png', dpi=300)
        chart_bytes.seek(0)
        plt.close()
        has_plot = True
    except Exception as chart_err:
        print(f"Chart render warning: {chart_err}")

    # Build PPTX File
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = params.get('title', 'Data Engine Report')
    slide.placeholders[1].text = f"Prepared for: {params.get('persona', 'Management')}"
    
    # Bullet Content Slides
    bullet_layout = prs.slide_layouts[1]
    for slide_data in deck_data.slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_data.title
        tf = slide.placeholders[1].text_frame
        for i, bullet in enumerate(slide_data.bullet_points):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                
    # Insert Generated Graph Slide
    if has_plot:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Visualized Insights: {chart_spec.chart_title}"
        slide.shapes.add_picture(chart_bytes, Inches(1.0), Inches(2.0), width=Inches(8.0))
        
    output_stream = io.BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    
    st.session_state["final_pptx_bytes"] = output_stream.getvalue()
    
    return {
        "output_ready": True,
        "messages": [AIMessage(content="✅ **Success!** Your presentation has been fully compiled with the embedded data visualization.")]
    }

def route_step(state: PresentationState):
    if len(state.get("missing_params", [])) > 0:
        return "gather_parameters"
    if state["messages"][-1].content.strip().lower() == "approve":
        return "generate_presentation_content"
    return END

# ==========================================
# 5. Persistent Memory & Engine Compilation
# ==========================================
@st.cache_resource
def compile_engine_workflow():
    builder = StateGraph(PresentationState)
    builder.add_node("ingest_and_summarize", ingest_and_summarize)
    builder.add_node("gather_parameters", gather_parameters)
    builder.add_node("generate_presentation_content", generate_presentation_content)

    builder.add_edge(START, "ingest_and_summarize")
    builder.add_edge("ingest_and_summarize", "gather_parameters")
    builder.add_conditional_edges(
        "gather_parameters", route_step,
        {"gather_parameters": END, "generate_presentation_content": "generate_presentation_content", END: END}
    )
    builder.add_edge("generate_presentation_content", END)
    return builder.compile(checkpointer=MemorySaver())

workflow_engine = compile_engine_workflow()

# ==========================================
# 6. Streamlit User Interface Layout
# ==========================================
st.title("📊 Multi-Horizon Data Profiler & Presentation Generator")

if "thread_id" not in st.session_state:
    st.session_state["thread_id"] = "persistent_session_final"
if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = []
if "file_processed" not in st.session_state:
    st.session_state["file_processed"] = False
if "final_pptx_bytes" not in st.session_state:
    st.session_state["final_pptx_bytes"] = None

config = {"configurable": {"thread_id": st.session_state["thread_id"]}}

# --- Sidebar ---
with st.sidebar:
    st.header("🗂️ Data Control Panel")
    uploaded_file = st.file_uploader("Upload Data Matrix (CSV or XLSX)", type=["csv", "xlsx"])
    
    if uploaded_file and not st.session_state["file_processed"]:
        with st.spinner("Processing file schema..."):
            try:
                if uploaded_file.name.endswith(".csv"):
                    raw_df = pd.read_csv(uploaded_file)
                else:
                    raw_df = pd.read_excel(uploaded_file)
                
                st.session_state["cached_df"] = raw_df
                schema_str = f"Columns: {raw_df.columns.tolist()}\nData Types: {raw_df.dtypes.to_dict()}\nDimensions: {raw_df.shape}"
                
                initial_state = {
                    "schema_info": schema_str,
                    "dataframe_json": raw_df.to_json(),
                    "missing_params": ["commentary_type", "persona", "topics", "pages", "title"],
                    "presentation_params": {},
                    "output_ready": False,
                    "messages": []
                }
                
                for event in workflow_engine.stream(initial_state, config):
                    for value in event.values():
                        if "messages" in value:
                            st.session_state["chat_history"].append(("assistant", value["messages"][-1].content))
                            
                st.session_state["file_processed"] = True
                st.rerun()
            except Exception as e:
                st.error(f"Error handling file parsing: {e}")

    if st.button("Reset Workspace Session", use_container_width=True):
        st.session_state.clear()
        st.rerun()

# --- Main Application Interface ---
if st.session_state["file_processed"]:
    df = st.session_state["cached_df"]
    
    with st.expander("🔍 Preview Uploaded Raw Data Matrix", expanded=False):
        st.dataframe(df, use_container_width=True)
        st.caption(f"Showing sample metrics — Total Rows: {df.shape[0]} | Total Columns: {df.shape[1]}")
        
    tab1, tab2 = st.tabs(["📉 Data Quality & Variance Analysis", "💬 Commentary & Presentation Builder"])
    
    # PAGE 1: Data Profiling
    with tab1:
        st.subheader("Data Quality Audit Summary")
        col1, col2, col3 = st.columns(3)
        
        missing_count = df.isnull().sum().sum()
        col1.metric("Missing Target Values", missing_count, delta=f"{'Action Needed' if missing_count > 0 else 'Clean Data'}")
        col2.metric("Duplicate Row Footprints", df.duplicated().sum())
        
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        col3.metric("Numeric Features Quantified", len(numeric_cols))
            
        st.markdown("---")
        st.subheader("Variance & Trend Aggregation")
        if len(numeric_cols) >= 1:
            selected_var_col = st.selectbox("Select Target Numeric Feature for Variance Tracking:", numeric_cols)
            v_col1, v_col2 = st.columns([1, 2])
            with v_col1:
                st.write("**Statistical Parameters:**")
                st.dataframe(df[selected_var_col].describe(), use_container_width=True)
            with v_col2:
                fig, ax = plt.subplots(figsize=(6, 3.5))
                sns.histplot(df[selected_var_col], kde=True, ax=ax, color="#1f77b4")
                ax.set_title(f"Distribution Variance: {selected_var_col}")
                st.pyplot(fig)
        else:
            st.info("No numeric column fields found.")

    # PAGE 2: Agentic Chat & Generation UI
    with tab2:
        st.subheader("Interactive Presentation Builder")
        
        chat_container = st.container(height=400)
        with chat_container:
            for role, message in st.session_state["chat_history"]:
                with st.chat_message(role):
                    st.markdown(message)
                    
        # Extract graph state memory
        state_snapshot = workflow_engine.get_state(config)
        state_values = state_snapshot.values if state_snapshot else {}
        missing_params = state_values.get("missing_params", [])
        output_ready = state_values.get("output_ready", False)

        st.markdown("---")
        
        # --- CONDITIONAL UI RENDERING ---
        if not output_ready:
            if missing_params:
                current_param = missing_params[0]
                suggestion_map = {
                    "commentary_type": ["Executive Summary", "Deep Dive Analysis", "Financial Overview"],
                    "persona": ["C-Suite Executives", "Technical Team", "General Stakeholders"],
                    "topics": ["Revenue & Margins", "Risk & Anomalies", "Cost Optimization"],
                    "pages": ["5", "8", "12"],
                    "title": ["Automated Data Insights", "Quarterly Performance", "Strategic Business Analysis"]
                }
                options = suggestion_map.get(current_param, ["Option 1", "Option 2", "Option 3"])
                
                chosen_suggestion = st.radio(
                    f"💡 Suggestions for **{current_param.replace('_', ' ').title()}**:", 
                    options, 
                    horizontal=True
                )
                
                with st.form("input_form", clear_on_submit=True):
                    user_text = st.text_input("✍️ Edit suggestion or type custom input:", value=chosen_suggestion)
                    submit = st.form_submit_button("Send 📤")
            else:
                with st.form("approve_form", clear_on_submit=True):
                    user_text = st.text_input("Final Step:", value="Approve")
                    submit = st.form_submit_button("Generate Deck 🚀")

            if submit and user_text:
                st.session_state["chat_history"].append(("user", user_text))
                with st.spinner("Processing next step..."):
                    for event in workflow_engine.stream({"messages": [HumanMessage(content=user_text)]}, config):
                        for value in event.values():
                            if "messages" in value:
                                st.session_state["chat_history"].append(("assistant", value["messages"][-1].content))
                st.rerun()
                
        else:
            # Output is ready: Display Download Button Inline
            st.success("🎉 Presentation compiled successfully! The deck includes your structured insights and the automated graph.")
            
            final_title = state_values.get('presentation_params', {}).get('title', 'Automated_Insights')
            safe_title = final_title.replace(" ", "_")
            
            st.download_button(
                label="📥 Download Generated Presentation (.pptx)",
                data=st.session_state["final_pptx_bytes"],
                file_name=f"{safe_title}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary"
            )
else:
    st.info("👋 System Idle. Please upload an Excel or CSV file in the sidebar to populate the dashboard modules.")
