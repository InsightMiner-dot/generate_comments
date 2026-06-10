import os
import io
import json
import base64
import pandas as pd
import matplotlib.pyplot as plt
from typing import TypedDict, Annotated, Dict, List, Any, Optional
from pydantic import BaseModel, Field

from langgraph.graph import StateGraph, START, END
from langgraph.graph.message import add_messages
from langgraph.checkpoint.memory import MemorySaver
from langchain_core.messages import AIMessage, SystemMessage, HumanMessage
from langchain_openai import AzureChatOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- 1. LLM Setup ---
azure_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
azure_api_key = os.getenv("AZURE_OPENAI_KEY")
azure_deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT")
azure_api_version = os.getenv("AZURE_OPENAI_API_VERSION")
model_name = os.getenv("MODEL_NAME")

llm = AzureChatOpenAI(
    azure_endpoint=azure_endpoint,
    api_key=azure_api_key,
    api_version=azure_api_version,
    deployment_name=azure_deployment,
    model=model_name,
    temperature=0.1
)

# --- 2. Advanced Presentation Content Schemas ---
class SuggestionList(BaseModel):
    suggestions: List[str] = Field(description="Exactly 3 short, contextually valid choices.")

class VarianceRow(BaseModel):
    metric_name: str = Field(description="Name of the parameter being evaluated")
    baseline: str = Field(description="Budget, baseline, or targeted baseline figure")
    actual: str = Field(description="The actual recorded data milestone figure")
    deviation: str = Field(description="Calculated absolute or directional variance delta")
    percentage: str = Field(description="Calculated percentage variance deviation (e.g. +14.2%)")

class RootCauseChain(BaseModel):
    symptom: str = Field(description="The core data anomaly or observed performance drop")
    contributing_factors: List[str] = Field(description="2-3 driving root elements or statistical vectors uncovered")
    remediation: str = Field(description="The targeted action plan or strategic resolution directive")

class SlideContent(BaseModel):
    slide_index: int = Field(description="Sequential index starting at 1")
    title: str = Field(description="Slide title capturing a concrete insight or data finding.")
    layout_type: str = Field(description="Must be one of: 'bullet_layout', 'table_grid', 'variance_analysis', 'root_cause'")
    bullet_points: List[str] = Field(description="3-4 impactful bullets containing real figures. Used for regular bullet layouts.")
    table_headers: List[str] = Field(default=[], description="Populated ONLY if layout_type is 'table_grid'")
    table_rows: List[List[str]] = Field(default=[], description="Populated cells matrix matching headers ONLY if layout_type is 'table_grid'")
    variance_data: List[VarianceRow] = Field(default=[], description="Populated ONLY if layout_type is 'variance_analysis'")
    root_cause_data: Optional[RootCauseChain] = Field(default=None, description="Populated ONLY if layout_type is 'root_cause'")

class PresentationDeck(BaseModel):
    slides: List[SlideContent]

class ChartSpecification(BaseModel):
    chart_title: str = Field(description="Title of the chart")
    x_column: str = Field(description="Exact column name for X-axis from dataset")
    y_column: str = Field(description="Exact column name for Y-axis")
    chart_type: str = Field(description="'bar', 'line', 'scatter', 'histogram', 'boxplot'")
    key_takeaways: List[str] = Field(description="3 concrete bullet points explaining what this data trend means.")

class PresentationState(TypedDict):
    messages: Annotated[list, add_messages]
    schema_info: str
    data_summary: str
    dataframe_json: str  
    current_step: str 
    suggestions: List[str]
    persona: str
    topics: str
    pages: str
    title: str
    graph_request: str
    draft_slides: List[Dict[str, Any]]
    output_ready: bool
    final_pptx_bytes: bytes
    chart_img_base64: str

# --- Premium Executive Dark Theme Palette Constants ---
DARK_BG = (15, 23, 42)          # #0F172A - Deep Dark Slate Background Canvas
TEXT_LIGHT = (248, 250, 252)    # #F8FAFC - High-Contrast Off-White Header Font
ACCENT_BLUE = (56, 189, 248)    # #38BDF8 - Vibrant Sky Blue Accent Typography
MUTED_TEXT = (148, 163, 184)    # #94A3B8 - Highly Readable Slate Gray body text
CRIT_RED = (239, 68, 68)        # #EF4444 - Contrast Deviation Alert Indicator
TABLE_CELL_BG = (241, 245, 249) # #F1F5F9 - Ultra Light Gray Table Grid Backgrounds
TEXT_DARK = (0, 0, 0)           # #000000 - High-Contrast Pure Black Grid Cell Content

def add_styled_text(tf, text, size_pt, bold=False, color_rgb=TEXT_LIGHT, is_first=False, space_after=12):
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color_rgb)
    p.space_after = Pt(space_after)
    return p

def apply_dark_background(slide):
    left = top = Inches(0)
    width = Inches(13.333)
    height = Inches(7.5)
    bg_shape = slide.shapes.add_shape(1, left, top, width, height) # 1 maps to Rectangle Geometry
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(*DARK_BG)
    bg_shape.line.fill.background()

# --- Chart Generation Helper ---
def generate_chart_img_base64(state: PresentationState, graph_request_str: str) -> tuple[str, Any]:
    if not graph_request_str or "no graph" in graph_request_str.lower():
        return "", None
    try:
        graph_prompt = f"Schema: {state['schema_info']}\nUser requested graph: {graph_request_str}\nDetermine chart config."
        chart_spec = llm.with_structured_output(ChartSpecification).invoke([SystemMessage(content=graph_prompt)])
        df = pd.read_json(io.StringIO(state["dataframe_json"]))
        
        plt.rcParams['font.family'] = 'sans-serif'
        plt.rcParams['font.sans-serif'] = ['DejaVu Sans', 'Segoe UI']
        
        fig, ax = plt.subplots(figsize=(7, 4.5), dpi=300, facecolor='#0F172A')
        ax.set_facecolor('#0F172A')
        
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#334155')
        ax.spines['bottom'].set_color('#334155')
        ax.set_axisbelow(True)
        ax.yaxis.grid(True, color='#1E293B', linestyle='-', linewidth=1)
        
        x_col = chart_spec.x_column if chart_spec.x_column in df.columns else df.columns[0]
        y_col = chart_spec.y_column if chart_spec.y_column in df.columns else df.select_dtypes(include=['number']).columns[0]
        
        sample_df = df.head(12).copy()
        sample_df[x_col] = sample_df[x_col].astype(str)
        c_type = chart_spec.chart_type.lower()
        
        if 'histogram' in c_type:
            ax.hist(df[y_col].dropna(), bins=10, color="#38bdf8", edgecolor='#0F172A')
        elif 'boxplot' in c_type:
            ax.boxplot(sample_df[y_col].dropna(), patch_artist=True, boxprops=dict(facecolor="#1e293b", color="#38bdf8"), medianprops=dict(color="#ef4444"))
        elif 'scatter' in c_type:
            ax.scatter(sample_df[x_col], sample_df[y_col], color="#ef4444", s=70, alpha=0.9, edgecolors='#f8fafc')
        else:
            bars = ax.bar(sample_df[x_col], sample_df[y_col], color="#38bdf8", width=0.55)
            for bar in bars:
                yval = bar.get_height()
                if yval > 0: ax.text(bar.get_x() + bar.get_width()/2, yval + (yval * 0.01), f"{yval:,.0f}", ha='center', va='bottom', fontsize=7, color='#94a3b8')
            
        ax.set_title(chart_spec.chart_title, fontweight='bold', fontsize=11, pad=15, color="#F8FAFC", loc='left')
        ax.tick_params(colors='#94a3b8', labelsize=8)
        plt.xticks(rotation=35, ha='right')
        plt.tight_layout()
        
        chart_bytes = io.BytesIO()
        plt.savefig(chart_bytes, format='png', dpi=300, facecolor=fig.get_facecolor(), edgecolor='none')
        chart_bytes.seek(0)
        base64_str = base64.b64encode(chart_bytes.getvalue()).decode('utf-8')
        plt.close()
        return base64_str, chart_spec
    except Exception as e:
        print(f"Plot Preview Error: {e}")
        return "", None

# --- Graph Orchestration Nodes ---

def ingest_and_summarize(state: PresentationState):
    """Step 0: Evaluates data profile and prepares system orchestrator states."""
    msg = ("### 📊 Data Profile Loaded & Synchronized\n"
           "Your dataset metrics are cataloged and visible in your **Data Quality Lab Dashboard**. "
           "You can use the pipeline toolkit on the left to treat anomalies, or proceed straight to architecture design.\n\n"
           "**Who is the target persona or audience for this presentation deck?**")
    return {
        "current_step": "persona",
        "suggestions": ["Data Analyst (The Gatekeeper)", "Financial Analyst (The Monetizer)", "Marketing Team (The Behavioralist)", "Executive Board (The Strategist)"],
        "messages": [AIMessage(content=msg)]
    }

def process_wizard_step(state: PresentationState):
    step = state.get("current_step")
    user_input = state["messages"][-1].content.strip()
    sugg_llm = llm.with_structured_output(SuggestionList)

    if step == "done" and any(w in user_input.lower() for w in ["yes", "restart", "new"]):
        return {
            "current_step": "persona", "draft_slides": [], "output_ready": False, "chart_img_base64": "",
            "suggestions": ["Data Analyst (The Gatekeeper)", "Financial Analyst (The Monetizer)", "Marketing Team (The Behavioralist)", "Executive Board (The Strategist)"],
            "messages": [AIMessage(content="🔄 **Session Reset!**\n\n**Who is the target audience or persona for this new deck?**")]
        }

    if step == "persona":
        prompt = f"Data Summary Overview Metrics:\n{state['data_summary']}\nPersona Selected: {user_input}\nSuggest 3 highly distinct scenario topics matching this persona."
        topics = sugg_llm.invoke([SystemMessage(content=prompt)]).suggestions
        return {
            "persona": user_input, "current_step": "topics", "suggestions": topics,
            "messages": [AIMessage(content=f"Got it. Target Audience: **{user_input}**.\n\nWhat specific **topics or operational scenarios** should we focus on?")]
        }

    elif step == "topics":
        return {
            "topics": user_input, "current_step": "pages", "suggestions": ["3 Slides", "5 Slides", "7 Slides"],
            "messages": [AIMessage(content="Excellent. How many **pages/slides** do you want to target inside this deck?")]
        }

    elif step == "pages":
        prompt = f"Dataset Metrics Profile Summary Context: {state['data_summary']}\nPersona Group: {state['persona']}\nFocus Area: {state['topics']}\nSuggest 3 executive, metric-focused Presentation Deck Titles."
        titles = sugg_llm.invoke([SystemMessage(content=prompt)]).suggestions
        return {
            "pages": user_input, "current_step": "title", "suggestions": titles,
            "messages": [AIMessage(content=f"Noted. What should be the **Main Title** of this PowerPoint presentation?")]
        }

    elif step == "title":
        p_lower = str(state.get("persona", "")).lower()
        if "data analyst" in p_lower:
            persona_pref = "Data Analyst (The Gatekeeper): Focuses on distributions, anomalies. Prefers Histograms, Box Plots, Scatter Plots."
        elif "financial" in p_lower:
            persona_pref = "Financial Analyst (The Monetizer): Focuses on profit margins, costs, budgets variances. Prefers Bar/Line charts."
        elif "marketing" in p_lower:
            persona_pref = "Marketing Team (The Behavioralist): Focuses on segmentation matrices, engagement trends. Prefers Funnel/Bar charts."
        else:
            persona_pref = "Executive Board (The Strategist): Focuses on high-level macro trajectories, strategy KPI health aggregates. Prefers trend lines or summary bars."

        prompt = f"Dataset Columns Schema: {state['schema_info']}\nAudience Guidelines Focus: {persona_pref}\nSuggest 3 contextually optimized visualizations matplotlib can draw based on fields."
        graphs = sugg_llm.invoke([SystemMessage(content=prompt)]).suggestions
        return {
            "title": user_input, "current_step": "graph", "suggestions": graphs + ["No Graph Needed"],
            "messages": [AIMessage(content=f"Great title choices! Based on your target persona and schema columns, what **graph or visualization** should we embed?")]
        }

    elif step == "graph":
        p_lower = str(state.get("persona", "")).lower()
        if "data analyst" in p_lower:
            audience_guardrail = "AUDIENCE: DATA ANALYST (THE GATEKEEPER). Focus heavily on statistical distributions, variance limits, anomalies, and mathematical structural verification metrics."
        elif "financial" in p_lower:
            audience_guardrail = "AUDIENCE: FINANCIAL ANALYST (THE MONETIZER). Focus heavily on operating profit margins, variances vs budgets, financial risks, and capital performance outcomes."
        elif "marketing" in p_lower:
            audience_guardrail = "AUDIENCE: MARKETING TEAM (THE BEHAVIORALIST). Focus heavily on customer behavior segments, retention vectors, funnel conversions, and ROI metrics."
        else:
            audience_guardrail = "AUDIENCE: EXECUTIVE BOARD (THE STRATEGIST). Focus completely on macro trajectories, strategic objectives, high-level corporate trajectories, and aggregated business health indicators."

        chart_b64, _ = generate_chart_img_base64(state, user_input)

        generation_prompt = f"""
        You are an elite corporate slide deck content architect. Create a comprehensive presentation blueprint framework outline.
        {audience_guardrail}
        Data Summary Profile Foundation Context: {state['data_summary']}
        Blueprint Directives:
        - Title: {state['title']} | Audience: {state['persona']} | Core Objective Focus: {state['topics']} | Slide Count: {state['pages']} Slides
        
        CRITICAL ENGINE STRUCTURE RULES:
        - If presenting clean bulleted textual lists, use layout_type 'bullet_layout'.
        - If displaying raw multi-column tables, matrices, or row metrics datasets summaries, use layout_type 'table_grid' and fill table_headers and table_rows data blocks.
        - If demonstrating budgets vs actual performance metrics comparisons, you MUST use layout_type 'variance_analysis' and fill variance_data rows.
        - If tracing down an error trend or performance breakdown anomaly, you MUST use layout_type 'root_cause' and populate the root_cause_data fields.
        """
        deck_draft = llm.with_structured_output(PresentationDeck).invoke([SystemMessage(content=generation_prompt)])
        slides_dict = [slide.model_dump() for slide in deck_draft.slides]
        
        msg = ("### 🛠️ Presentation Strategy Structural Outline Compiled\n"
               "The slide blueprint design framework and visualization preview are processed. "
               "You can **manually edit any slide card title, bullet text array, or table cell value on the right panel**, "
               "or ask me to apply transformations here. Click **Approve Plan & Compile Presentation** to generate your PowerPoint deck.")
        
        return {
            "graph_request": user_input, "current_step": "review_slides", "draft_slides": slides_dict, "chart_img_base64": chart_b64,
            "suggestions": ["Approve Plan & Compile Presentation", "Make headings punchier"], "messages": [AIMessage(content=msg)]
        }

    elif step == "review_slides":
        if "approve" in user_input.lower() or "compile" in user_input.lower():
            return {"current_step": "generating", "suggestions": [], "messages": [AIMessage(content="⚙️ Layout approved. Executing multi-agent validation loops and writing widescreen binary streams...")]}
        else:
            edit_prompt = f"""Modify this existing slide blueprint layout matching instructions. Keep data observations structurally sound.
            Data Metrics Context: {state['data_summary']}
            Current Slide Structure Layout: {json.dumps(state.get('draft_slides', []))}
            User Mutation Command: "{user_input}" """
            updated_deck = llm.with_structured_output(PresentationDeck).invoke([SystemMessage(content=edit_prompt)])
            slides_dict = [slide.model_dump() for slide in updated_deck.slides]
            return {
                "draft_slides": slides_dict, "suggestions": ["Approve Plan & Compile Presentation"],
                "messages": [AIMessage(content=f"🔄 Blueprint updated with request: *\"{user_input}\"*. Review changes in the preview panel.")]
            }

def generate_presentation(state: PresentationState):
    slides_list = list(state.get("draft_slides", []))
    audit_prompt = f"""
    You are a meticulous Multi-Agent Data Compliance Inspector. Cross-examine the text inside this slide blueprint against the true dataset statistics.
    Slide Blueprint Frame: {json.dumps(slides_list)}
    True Dataset Statistics Snapshot: {state['data_summary']}
    
    If any calculation, total, percentage, or milestone listed on the slides does not match the true summary data, you MUST modify the values to align them. Ensure 100% numerical validation.
    """
    audited_deck = llm.with_structured_output(PresentationDeck).invoke([SystemMessage(content=audit_prompt)])
    validated_slides = [slide.model_dump() for slide in audited_deck.slides]

    chart_b64, chart_spec = generate_chart_img_base64(state, state["graph_request"])
    chart_bytes = io.BytesIO(base64.b64decode(chart_b64)) if chart_b64 else None

    # PowerPoint Compiler Construction Engine (16:9 widescreen config)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] 

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide)
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    add_styled_text(tf, state.get('title', 'Generated Report'), 44, bold=True, color_rgb=TEXT_LIGHT, is_first=True, space_after=18)
    add_styled_text(tf, f"Target Operational Focus: {state.get('persona', 'Enterprise Stakeholders')}", 20, bold=False, color_rgb=ACCENT_BLUE)

    # --- Framework Dispatcher Render Loop ---
    for raw_slide in validated_slides:
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_background(slide)
        
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        add_styled_text(tf_title, raw_slide.get("title", "Slide"), 32, bold=True, color_rgb=TEXT_LIGHT, is_first=True)
        
        l_type = raw_slide.get("layout_type")
        
        # FIXED CONTRAST HOOK 1: VARIANCE ANALYSIS TABLE COMPILER (Light Cell Fill, Pure Black Text)
        if l_type == "variance_analysis" and raw_slide.get("variance_data"):
            v_data = raw_slide["variance_data"]
            table_shape = slide.shapes.add_table(len(v_data) + 1, 5, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5)).table
            headers = ["Metric Parameter", "Target/Baseline", "Actual Performance", "Delta Variance", "Deviation %"]
            for c_idx, h_text in enumerate(headers):
                cell = table_shape.cell(0, c_idx)
                cell.text = h_text
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59) # Deep Blue Gray Header Row
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*ACCENT_BLUE)
                cell.text_frame.paragraphs[0].font.bold = True
                
            for r_idx, row in enumerate(v_data):
                for c_idx, val in enumerate([row["metric_name"], row["baseline"], row["actual"], row["deviation"]]):
                    cell = table_shape.cell(r_idx + 1, c_idx)
                    cell.text = str(val)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*TABLE_CELL_BG)
                    p = cell.text_frame.paragraphs[0]
                    p.font.color.rgb = RGBColor(*TEXT_DARK) # HIGH-CONTRAST BLACK CONTENT FONT
                    p.font.size = Pt(13)
                
                p_cell = table_shape.cell(r_idx + 1, 4)
                p_cell.text = str(row["percentage"])
                p_cell.fill.solid()
                p_cell.fill.fore_color.rgb = RGBColor(*TABLE_CELL_BG)
                p = p_cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(13)
                if "-" in str(row["percentage"]):
                    p.font.color.rgb = RGBColor(*CRIT_RED)
                else:
                    p.font.color.rgb = RGBColor(16, 185, 129)

        # RUNTIME FRAMEWORK 2: ROOT CAUSE BLOCK SCHEMATIC
        elif l_type == "root_cause" and raw_slide.get("root_cause_data"):
            rc = raw_slide["root_cause_data"]
            
            s_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.4), Inches(4.5))
            s_tf = s_box.text_frame
            s_tf.word_wrap = True
            add_styled_text(s_tf, "🚨 Observed Symptom Anomaly", 14, bold=True, color_rgb=CRIT_RED, is_first=True)
            add_styled_text(s_tf, rc["symptom"], 15, bold=False, color_rgb=TEXT_LIGHT)
            
            c_box = slide.shapes.add_textbox(Inches(4.9), Inches(2.0), Inches(3.8), Inches(4.5))
            c_tf = c_box.text_frame
            c_tf.word_wrap = True
            add_styled_text(c_tf, "🔍 Core Driving Vectors", 14, bold=True, color_rgb=ACCENT_BLUE, is_first=True)
            for factor in rc["contributing_factors"]:
                add_styled_text(c_tf, f"• {factor}", 15, bold=False, color_rgb=MUTED_TEXT, space_after=8)
                
            r_box = slide.shapes.add_textbox(Inches(9.2), Inches(2.0), Inches(3.1), Inches(4.5))
            r_tf = r_box.text_frame
            r_tf.word_wrap = True
            add_styled_text(r_tf, "⚡ Remediation Action Plan", 14, bold=True, color_rgb=(16, 185, 129), is_first=True)
            add_styled_text(r_tf, rc["remediation"], 15, bold=False, color_rgb=TEXT_LIGHT)

        # FIXED CONTRAST HOOK 2: DATA TABLE GRID MATRIX COMPILER (Light Cell Fill, Pure Black Text)
        elif l_type == "table_grid" and raw_slide.get("table_headers"):
            headers = raw_slide["table_headers"]
            rows = raw_slide["table_rows"]
            table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5)).table
            for c_idx, h_text in enumerate(headers):
                cell = table_shape.cell(0, c_idx)
                cell.text = str(h_text)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*ACCENT_BLUE)
                cell.text_frame.paragraphs[0].font.bold = True
                
            for r_idx, row_data in enumerate(rows):
                for c_idx, cell_val in enumerate(row_data):
                    cell = table_shape.cell(r_idx + 1, c_idx)
                    cell.text = str(cell_val)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*TABLE_CELL_BG)
                    p = cell.text_frame.paragraphs[0]
                    p.font.color.rgb = RGBColor(*TEXT_DARK) # HIGH-CONTRAST BLACK CONTENT FONT
                    p.font.size = Pt(13)
        else:
            # Layout Type 4: Classic corporate widescreen bullet templates
            body_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0))
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            for idx, bullet in enumerate(raw_slide.get("bullet_points", [])):
                add_styled_text(tf_body, f"•  {bullet}", 16, bold=False, color_rgb=MUTED_TEXT, is_first=(idx == 0), space_after=14)
                
    # --- Side-By-Side Visual Metric Takeaways slide ---
    if chart_bytes and chart_spec:
        slide = prs.slides.add_slide(blank_layout)
        apply_dark_background(slide)
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
        add_styled_text(title_box.text_frame, chart_spec.chart_title, 32, bold=True, color_rgb=TEXT_LIGHT, is_first=True)
        
        slide.shapes.add_picture(chart_bytes, Inches(0.6), Inches(1.8), width=Inches(6.8))
        explanation_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.8), Inches(4.8))
        tf_explain = explanation_box.text_frame
        tf_explain.word_wrap = True
        add_styled_text(tf_explain, "Analytical Takeaways", 20, bold=True, color_rgb=ACCENT_BLUE, is_first=True, space_after=14)
        for insight in chart_spec.key_takeaways:
            add_styled_text(tf_explain, f"• {insight}", 15, bold=False, color_rgb=MUTED_TEXT, is_first=False, space_after=12)
        
    output_stream = io.BytesIO()
    prs.save(output_stream)
    
    return {
        "output_ready": True, "current_step": "done", "final_pptx_bytes": output_stream.getvalue(),
        "suggestions": ["Yes, Start New Session"],
        "messages": [AIMessage(content="🎉 **Success!** Your presentation strategy has been audited and compiled into a premium executive high-contrast dark PowerPoint presentation deck.")]
    }

def start_router(state: PresentationState):
    if state.get("current_step") == "init": return "ingest_and_summarize"
    return "process_wizard_step"

def router(state: PresentationState):
    if state["current_step"] == "generating": return "generate_presentation"
    return "process_wizard_step"

workflow = StateGraph(PresentationState)
workflow.add_node("ingest_and_summarize", ingest_and_summarize)
workflow.add_node("process_wizard_step", process_wizard_step)
workflow.add_node("generate_presentation", generate_presentation)

workflow.add_conditional_edges(START, start_router, {"ingest_and_summarize": "ingest_and_summarize", "process_wizard_step": "process_wizard_step"})
workflow.add_edge("ingest_and_summarize", END)
workflow.add_conditional_edges("process_wizard_step", router, {"process_wizard_step": END, "generate_presentation": "generate_presentation"})
workflow.add_edge("generate_presentation", END)

memory = MemorySaver()
app_engine = workflow.compile(checkpointer=memory)
